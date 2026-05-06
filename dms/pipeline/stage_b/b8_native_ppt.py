"""
B-8: Native PowerPoint Processor (.pptx専用)

python-pptx を使用して、ネイティブPowerPointファイルから構造化データを抽出。
スライド・シェイプ単位でテキストを回収、座標順にソート。
"""

from pathlib import Path
from typing import Dict, Any, List
from loguru import logger


class B8NativePPTProcessor:
    """B-8: Native PowerPoint Processor (.pptx専用)"""

    def process(self, file_path: Path, log_file=None) -> Dict[str, Any]:
        """
        .pptx ファイルから構造化データを抽出

        Args:
            file_path: .pptxファイルパス
            log_file: 個別ログファイルパス（Noneなら共有ロガーのみ）

        Returns:
            {
                'is_structured': bool,
                'text_with_tags': str,           # 全スライドのテキスト
                'structured_tables': [...],      # 表構造データ
                'slides': [...],                 # スライド情報
                'tags': {...},                   # メタ情報
                'media_elements': [...]          # 埋め込み画像
            }
        """
        _sink_id = None
        if log_file:
            _sink_id = logger.add(
                str(log_file),
                format="{time:HH:mm:ss} | {level:<5} | {message}",
                filter=lambda r: "[B-8]" in r["message"],
                level="DEBUG",
                encoding="utf-8",
            )
        try:
            return self._process_impl(file_path)
        finally:
            if _sink_id is not None:
                logger.remove(_sink_id)

    def _process_impl(self, file_path: Path) -> Dict[str, Any]:
        logger.info(f"[B-8] Native PowerPoint処理開始: {file_path.name}")

        try:
            from pptx import Presentation
        except ImportError:
            logger.error("[B-8] python-pptx がインストールされていません")
            return self._error_result("python-pptx not installed")

        try:
            prs = Presentation(str(file_path))

            # 全スライドを抽出
            slides = self._extract_slides(prs)

            # テキストを生成
            text_with_tags = self._build_text(slides)

            # 表を抽出
            tables = self._extract_tables(slides)

            # メタ情報
            tags = {
                'slide_count': len(slides),
                'has_table': len(tables) > 0,
                'total_shapes': sum(s.get('shape_count', 0) for s in slides)
            }

            logger.info(f"[B-8] 抽出完了: スライド={len(slides)}, 表={len(tables)}")
            for slide in slides:
                for shape_idx, shape in enumerate(slide.get('shapes', [])):
                    logger.info(f"[B-8] slide{slide.get('slide_num')} shape{shape_idx}: {shape.get('text', '')}")

            return {
                'is_structured': True,
                'text_with_tags': text_with_tags,
                'structured_tables': tables,
                'slides': slides,
                'tags': tags,
                'media_elements': []  # TODO: 画像抽出
            }

        except Exception as e:
            logger.error(f"[B-8] 処理エラー: {e}", exc_info=True)
            return self._error_result(str(e))

    def _extract_slides(self, prs) -> List[Dict[str, Any]]:
        """
        全スライドを抽出

        Returns:
            [{
                'slide_num': int,
                'shapes': [...],
                'shape_count': int
            }, ...]
        """
        slides = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            shapes = []

            for shape in slide.shapes:
                # テキストを持つシェイプのみ抽出
                if not hasattr(shape, "text") or not shape.text.strip():
                    continue

                # 座標情報を取得
                top = shape.top if hasattr(shape, 'top') else 0
                left = shape.left if hasattr(shape, 'left') else 0

                shapes.append({
                    'text': shape.text,
                    'top': top,
                    'left': left,
                    'shape_type': shape.shape_type
                })

            # 座標順にソート（上から下、左から右）
            shapes.sort(key=lambda s: (s['top'], s['left']))

            slides.append({
                'slide_num': slide_num,
                'shapes': shapes,
                'shape_count': len(shapes)
            })

        return slides

    def _extract_tables(self, slides: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """
        スライド内の表を抽出

        Returns:
            [{
                'slide_num': int,
                'table_index': int,
                'data': [[...], [...], ...]
            }, ...]
        """
        # TODO: python-pptx の Table オブジェクトから抽出
        # 現在は未実装（将来の拡張ポイント）
        return []

    def _build_text(self, slides: List[Dict[str, Any]]) -> str:
        """
        全スライドのテキストを生成

        Returns:
            スライド番号とシェイプテキストを結合
        """
        result = []

        for slide in slides:
            result.append(f"\n[SLIDE {slide['slide_num']}]")
            for shape in slide['shapes']:
                result.append(shape['text'])
            result.append(f"[/SLIDE]\n")

        return "\n".join(result)

    def _error_result(self, error_message: str) -> Dict[str, Any]:
        """エラー結果を返す"""
        return {
            'is_structured': False,
            'error': error_message,
            'text_with_tags': '',
            'structured_tables': [],
            'slides': [],
            'tags': {},
            'media_elements': []
        }
