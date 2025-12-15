#【実行場所】: ターミナルまたはVS Code
#【対象ファイル】: 新規作成
#【ファイルパス】: core/processors/office.py
#【実行方法】: 以下のコードをファイルにコピー＆ペーストして保存してください。

"""
Office プロセッサ (テキスト抽出)

設計書: COMPLETE_IMPLEMENTATION_GUIDE_v3.md の 1.4節に基づき、Officeファイルからテキストを抽出する。
対応形式: DOCX, XLSX, PPTX
"""
from typing import Dict, Any
from pathlib import Path
from io import BytesIO
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from loguru import logger

class OfficeProcessor:
    """DOCX, XLSX, PPTXファイルからテキストを抽出するプロセッサ"""
    
    def __init__(self):
        # logger.info("Officeプロセッサ初期化完了")
        pass

    def extract_from_docx(self, file_path: str) -> Dict[str, Any]:
        """DOCXファイルから全文を抽出する（Windows WinError 32対策済み）"""
        file_path = Path(file_path)
        full_text = []

        try:
            # BytesIOを使用してファイルロックを回避
            with open(file_path, 'rb') as f:
                file_data = BytesIO(f.read())

            document = Document(file_data)
            for paragraph in document.paragraphs:
                full_text.append(paragraph.text)

            # テーブルのテキストも追加
            for table in document.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text)
                    full_text.append(" | ".join(row_text))

            content = "\n".join(full_text)

            return {"content": content, "metadata": {"file_type": "docx"}, "success": True}

        except Exception as e:
            # logger.error(f"DOCXテキスト抽出エラー ({file_path}): {e}")
            return {"content": "", "metadata": {"file_type": "docx", "error": str(e)}, "success": False, "error_message": str(e)}

    def extract_from_xlsx(self, file_path: str) -> Dict[str, Any]:
        """XLSXファイルから全シートのテキストを抽出する（Windows WinError 32対策済み）"""
        file_path = Path(file_path)
        full_text = []

        try:
            # BytesIOを使用してファイルロックを回避
            with open(file_path, 'rb') as f:
                file_data = BytesIO(f.read())

            workbook = load_workbook(file_data, data_only=True)
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                full_text.append(f"\n---SHEET: {sheet_name}---\n")

                for row in sheet.iter_rows():
                    row_data = []
                    for cell in row:
                        value = str(cell.value) if cell.value is not None else ""
                        row_data.append(value)
                    full_text.append(" | ".join(row_data))

            # ワークブックを明示的に閉じる
            workbook.close()

            content = "\n".join(full_text)

            return {"content": content, "metadata": {"file_type": "xlsx"}, "success": True}

        except Exception as e:
            # logger.error(f"XLSXテキスト抽出エラー ({file_path}): {e}")
            return {"content": "", "metadata": {"file_type": "xlsx", "error": str(e)}, "success": False, "error_message": str(e)}

    def extract_from_pptx(self, file_path: str) -> Dict[str, Any]:
        """PPTXファイルから全スライドのテキストを抽出する（Windows WinError 32対策済み）"""
        file_path = Path(file_path)
        full_text = []

        try:
            # BytesIOを使用してファイルロックを回避
            with open(file_path, 'rb') as f:
                file_data = BytesIO(f.read())

            presentation = Presentation(file_data)
            for i, slide in enumerate(presentation.slides):
                full_text.append(f"\n---SLIDE: {i+1}---")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text.append(shape.text)

            content = "\n".join(full_text)

            return {"content": content, "metadata": {"file_type": "pptx"}, "success": True}

        except Exception as e:
            # logger.error(f"PPTXテキスト抽出エラー ({file_path}): {e}")
            return {"content": "", "metadata": {"file_type": "pptx", "error": str(e)}, "success": False, "error_message": str(e)}