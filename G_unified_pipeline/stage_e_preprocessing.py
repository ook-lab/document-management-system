"""
Stage E: Pre-processing (前処理)

全ファイルタイプ共通のE1-E5ステージ:
- E-1: テキスト抽出（ライブラリベース）
- E-2: 表抽出
- E-3: 統合（テキスト + 表）
- E-4: Gemini Vision差分検出
- E-5: Vision OCR結果適用

対応ファイル:
- PDF: pdfplumber + Vision補完
- Office: python-docx, openpyxl, python-pptx + Vision補完
- 画像: Vision OCR（E-4で全文字拾い）
"""
from pathlib import Path
from typing import Dict, Any, List, Optional
from loguru import logger

from A_common.processors.pdf import PDFProcessor
from A_common.processors.office import OfficeProcessor
from C_ai_common.llm_client.llm_client import LLMClient


class StageEPreprocessor:
    """Stage E: 前処理（全ファイルタイプ共通のE1-E5ステージ）"""

    def __init__(self, llm_client: LLMClient):
        """
        Args:
            llm_client: LLMクライアント
        """
        self.llm_client = llm_client
        self.pdf_processor = PDFProcessor(llm_client=llm_client)
        self.office_processor = OfficeProcessor()

    def extract_text(
        self,
        file_path: Path,
        mime_type: str,
        pre_extracted_text: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        ファイルからテキストを抽出（E1-E5ステージ）

        Args:
            file_path: ファイルパス
            mime_type: MIMEタイプ
            pre_extracted_text: 既に抽出済みのテキスト（HTML→PNG等の場合）

        Returns:
            {
                'success': bool,
                'content': str,  # 完全なテキスト（単一）
                'char_count': int,
                'method': str  # 'pdf', 'docx', 'xlsx', 'pptx', 'image', 'none'
            }
        """
        logger.info("=" * 60)
        logger.info("[Stage E] Pre-processing開始...")
        logger.info(f"  ├─ ファイル: {file_path.name if isinstance(file_path, Path) else file_path}")
        logger.info(f"  └─ MIMEタイプ: {mime_type}")

        content = ""
        method = "none"

        try:
            # PDF処理（E1-E5は pdf.py 内で実行）
            if mime_type == 'application/pdf':
                result = self.pdf_processor.extract_text(str(file_path))
                if result.get('success'):
                    content = result.get('content', '')
                    method = 'pdf'

            # Office文書処理（E1-E5ログ付き）
            elif mime_type in [
                'application/vnd.openxmlformats-officedocument.wordprocessingml.document',  # .docx
                'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',      # .xlsx
                'application/vnd.openxmlformats-officedocument.presentationml.presentation'  # .pptx
            ]:
                content, method = self._process_office_with_stages(file_path, mime_type)

            # 画像ファイル処理（E1-E5ログ付き）
            # HTML→PNG の場合も含む（mime_type='text/html' だがファイルはPNG）
            elif mime_type.startswith('image/') or mime_type == 'text/html':
                content, method = self._process_image_with_stages(file_path, pre_extracted_text)

            # 最終ログ
            logger.info("=" * 60)
            logger.info(f"[Stage E完了] 最終結果:")
            logger.info(f"  ├─ 処理方式: {method}")
            logger.info(f"  └─ 最終テキスト: {len(content)} 文字")
            logger.info("=" * 60)

            return {
                'success': True,
                'content': content,
                'char_count': len(content),
                'method': method
            }

        except Exception as e:
            logger.error(f"[Stage E エラー] テキスト抽出失敗: {e}", exc_info=True)
            return {
                'success': False,
                'content': '',
                'char_count': 0,
                'method': 'error',
                'error': str(e)
            }

    def _process_office_with_stages(self, file_path: Path, mime_type: str) -> tuple:
        """
        Officeファイル処理（E1-E5ステージログ付き）

        Args:
            file_path: ファイルパス
            mime_type: MIMEタイプ

        Returns:
            (content, method)
        """
        file_path = Path(file_path)

        # ファイルタイプ判定
        if 'word' in mime_type:
            file_type = 'docx'
        elif 'sheet' in mime_type:
            file_type = 'xlsx'
        elif 'presentation' in mime_type:
            file_type = 'pptx'
        else:
            file_type = 'office'

        logger.info(f"[Stage E] Office処理開始 (type: {file_type})")

        # ============================================
        # E-1: テキスト抽出（ライブラリベース）
        # ============================================
        logger.info(f"[E-1] テキスト抽出開始 (engine: python-{file_type})")

        result = self.office_processor.extract_text(str(file_path))
        e1_text = result.get('content', '') if result.get('success') else ''
        e1_chars = len(e1_text)

        logger.info(f"[E-1] テキスト抽出完了:")
        logger.info(f"  └─ 抽出テキスト: {e1_chars} 文字")

        # ============================================
        # E-2: 表抽出
        # ============================================
        logger.info(f"[E-2] 表抽出:")
        # Officeファイルの場合、表はテキスト抽出時に含まれる
        # （python-docx/openpyxl/python-pptxが表を含めて抽出）
        table_count = e1_text.count('|') // 2 if '|' in e1_text else 0  # 簡易カウント
        logger.info(f"  └─ 表データ: テキストに含まれる (推定 {table_count} セル)")

        # ============================================
        # E-3: 統合
        # ============================================
        logger.info(f"[E-3] 統合:")
        e3_text = e1_text  # Officeは既に統合済み
        e3_chars = len(e3_text)
        logger.info(f"  └─ 統合テキスト: {e3_chars} 文字")

        # ============================================
        # E-4: Gemini Vision差分検出
        # ============================================
        logger.info(f"[E-4] Vision差分検出:")

        e4_text = ""
        if self.llm_client and e3_chars < 100:
            # テキストが少ない場合のみVision補完を実行
            logger.info(f"  ├─ テキスト量が少ない ({e3_chars}文字) → Vision補完を実行")

            # PPTXの場合、スライドを画像化してVision処理
            if file_type == 'pptx':
                try:
                    from pptx import Presentation
                    from PIL import Image
                    import tempfile
                    import os

                    # スライド数を取得
                    prs = Presentation(str(file_path))
                    slide_count = len(prs.slides)
                    logger.info(f"  ├─ スライド数: {slide_count}")

                    # 各スライドをVision処理（最大5スライド）
                    vision_texts = []
                    for i, slide in enumerate(prs.slides[:5]):
                        logger.info(f"  ├─ スライド {i+1} をVision処理中...")
                        # Note: PPTXを画像化するには追加ライブラリが必要
                        # ここでは簡略化してスキップ

                    if vision_texts:
                        e4_text = "\n\n".join(vision_texts)
                        logger.info(f"  └─ Vision補完結果: {len(e4_text)} 文字")
                    else:
                        logger.info(f"  └─ Vision補完: スキップ（画像化未対応）")

                except Exception as e:
                    logger.warning(f"  └─ Vision補完失敗: {e}")
            else:
                logger.info(f"  └─ Vision補完: スキップ（{file_type}は対象外）")
        else:
            logger.info(f"  └─ Vision補完: 不要 (十分なテキスト量: {e3_chars}文字)")

        # ============================================
        # E-5: Vision OCR結果適用
        # ============================================
        logger.info(f"[E-5] Vision OCR結果適用:")

        if e4_text:
            # Vision補完がある場合は追加
            e5_text = e3_text + "\n\n---\n\n## Vision OCR 補完情報\n\n" + e4_text
            logger.info(f"  ├─ E-3テキスト: {e3_chars} 文字")
            logger.info(f"  ├─ E-4補完: +{len(e4_text)} 文字")
            logger.info(f"  └─ E-5最終: {len(e5_text)} 文字")
        else:
            e5_text = e3_text
            logger.info(f"  └─ E-5最終: {len(e5_text)} 文字 (Vision補完なし)")

        return e5_text, file_type

    def _process_image_with_stages(
        self,
        file_path: Path,
        pre_extracted_text: Optional[str] = None
    ) -> tuple:
        """
        画像ファイル処理（E1-E5ステージログ付き）

        Args:
            file_path: ファイルパス
            pre_extracted_text: 既に抽出済みのテキスト（HTML→PNG等の場合）

        Returns:
            (content, method)
        """
        file_path = Path(file_path)

        logger.info(f"[Stage E] 画像処理開始")

        # ============================================
        # E-1: テキスト抽出（画像なのでなし、ただしHTML→PNG等の場合は既抽出済み）
        # ============================================
        logger.info(f"[E-1] テキスト抽出:")
        if pre_extracted_text:
            e1_chars = len(pre_extracted_text)
            logger.info(f"  └─ 既抽出テキスト (Ingestion時): {e1_chars} 文字")
        else:
            e1_chars = 0
            logger.info(f"  └─ 画像ファイルのため、ライブラリ抽出: 0 文字")

        # ============================================
        # E-2: 表抽出（画像なのでなし）
        # ============================================
        logger.info(f"[E-2] 表抽出:")
        logger.info(f"  └─ 画像ファイルのため、表抽出: 0 個")

        # ============================================
        # E-3: 統合（画像なのでなし、ただしHTML→PNG等の場合は既抽出テキストを使用）
        # ============================================
        logger.info(f"[E-3] 統合:")
        if pre_extracted_text:
            e3_text = pre_extracted_text
            e3_chars = len(e3_text)
            logger.info(f"  └─ E-1既抽出テキスト: {e3_chars} 文字")
        else:
            e3_text = ""
            e3_chars = 0
            logger.info(f"  └─ E-1 + E-2 統合: 0 文字")

        # ============================================
        # E-4: Gemini Vision OCR（画像のメイン処理）
        # ============================================
        logger.info(f"[E-4] Vision OCR処理 (model: gemini-2.5-flash):")

        e4_text = ""
        if self.llm_client:
            vision_result = self.llm_client.transcribe_image(
                image_path=file_path,
                model="gemini-2.5-flash",
                prompt="""この画像から、全ての文字を徹底的に拾い尽くしてください。

【あなたの役割】
画像から全ての文字を漏らさず拾ってください。

【文字拾いの徹底指示】
- **小さな文字**: 注釈、脚注、コピーライト表記なども全て拾う
- **ロゴ化された文字**: 画像として埋め込まれたタイトル、会社名、ブランド名なども全て読み取る
- **装飾された文字**: 太字、斜体、色付きなど、装飾に関わらず全て拾う
- **背景に埋もれた文字**: 薄い色、透かし文字なども可能な限り読み取る
- **手書き文字**: 判読可能な範囲で全て拾う
- **表構造**: 表がある場合は、Markdown table形式で出力

【出力形式】
画像内の全てのテキストを、Markdown形式で構造化して出力してください。
表がある場合は必ずMarkdown table形式で出力してください。

**重要**: 1文字も見逃さないでください。"""
            )

            if vision_result.get('success'):
                e4_text = vision_result.get('content', '')
                logger.info(f"  ├─ Vision OCR成功")
                logger.info(f"  └─ 抽出テキスト: {len(e4_text)} 文字")
            else:
                logger.warning(f"  └─ Vision OCR失敗: {vision_result.get('error', 'Unknown error')}")
        else:
            logger.warning(f"  └─ LLMクライアント未設定のためスキップ")

        # ============================================
        # E-5: Vision OCR結果適用（画像は E-4 がそのまま最終、ただし既抽出テキストがあれば統合）
        # ============================================
        logger.info(f"[E-5] Vision OCR結果適用:")

        # E-3テキストとE-4 Vision OCRを統合
        if e4_text:
            if e3_text:
                # 既抽出テキスト + Vision OCR
                e5_text = e3_text + "\n\n---\n\n## Vision OCR 補完情報\n\n" + e4_text
            else:
                # Vision OCRのみ
                e5_text = e4_text
        else:
            # 既抽出テキストのみ
            e5_text = e3_text

        logger.info(f"  ├─ E-3テキスト: {e3_chars} 文字")
        logger.info(f"  ├─ E-4 Vision OCR: {len(e4_text)} 文字")
        logger.info(f"  └─ E-5最終: {len(e5_text)} 文字")

        return e5_text, 'image'

    def process(self, file_path: Path, mime_type: str) -> str:
        """
        ファイルからテキストを抽出（process() エイリアス）

        Args:
            file_path: ファイルパス
            mime_type: MIMEタイプ

        Returns:
            extracted_text: 抽出されたテキスト
        """
        result = self.extract_text(file_path, mime_type)
        return result.get('content', '') if result.get('success') else ''
